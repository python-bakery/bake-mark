import argparse
import json
import math
import time
import os
from pathlib import Path
import io

# Progress bar
from tqdm import tqdm

# Hashing stuff
from friendly_hash import hash, hash_exists

# Markdown parsing stuff
import marko
from marko.ext.gfm import GFMRendererMixin
from marko.ext.gfm import elements
import marko.renderer
from markdown_tools import extract_front_matter

# Code highlighting stuff
from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import html
from pygments.formatter import Formatter
from pygments.formatters.img import ImageFormatter
from pygments.formatters.html import HtmlFormatter
from pygments.util import ClassNotFound
from pygments.style import Style
from pygments import token
#from pygments.token import Token, Comment, Keyword, Name, String, \
#    Error, Generic, Number, Operator, Whitespace
from pygments.token import (
    Keyword,
    Name,
    Comment,
    String,
    Error,
    Number,
    Operator,
    Generic,
    Literal,
)
from pygments.styles.sas import SasStyle
from pygments.styles import get_style_by_name

# Powerpoint stuff
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from pptx.opc.package import PartFactory
from pptx.parts.media import MediaPart

for aud_type in [ 'audio/mp3', 'audio/mp4', 'audio/mid', 
                    'audio/x-wav', 'audio/mpeg' ]:
    PartFactory.part_type_for.update({aud_type: MediaPart})

# XML Handling stuff
from lxml import etree

# Amazon Polly stuff
import polly

# MP3 Handling stuff
from mutagen.mp3 import MP3

# Monkey Patches
import python_pptx_patches

# Windows communication client
import win32com.client

# Subtitling
from make_subtitles import make_captions

# FFMPEG conversion
import ffmpeg

# Local important data
from locations import POWERPOINT_TEMPLATE

# Actual code!

def replace_with_image(img, shape, slide, max_size=False, presentation=None):
    pic = slide.shapes.add_picture(img, shape.left, shape.top)

    # calculate max width/height for target size
    ratio = min(shape.width / float(pic.width), shape.height / float(pic.height))

    if max_size:
        start_of_content_area = slide.shapes.title.top + slide.shapes.title.height
        height_of_content_area = presentation.slide_height - start_of_content_area
        height_of_content_area -= Inches(.5)
        pic.width = int(height_of_content_area * pic.width / pic.height)
        pic.height = int(height_of_content_area)
        pic.top = start_of_content_area
        pic.left = shape.left
    else:
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        pic.top = shape.top + ((shape.height - pic.height) // 2)
        pic.left = shape.left + ((shape.width - pic.width) // 2)

    placeholder = shape.element
    placeholder.getparent().remove(placeholder)
    return


def no_bullet(paragraph):
    paragraph._pPr.insert(
        0,
        etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"),
    )


def _parse_extras(line):
    if not line:
        return {}
    return {k: json.loads(v)
            for part in line.split(",")
            for k, v in [part.split("=")]}


class CodeStyle(Style):
    default_style = ''
    styles = {
        token.Whitespace:            '#bbbbbb',
        token.Comment:               '#008800',
        token.String:                '#800080',
        token.Number:                '#2c8553',
        token.Other:                 'bg:#ffffe0',
        token.Keyword:               '#2c2cff',
        token.Keyword.Reserved:      '#353580',
        token.Keyword.Constant:      '',
        token.Name.Builtin:          '#2c2cff',
        token.Name.Variable:         '#2c2cff',
        token.Generic:               '#2c2cff',
        token.Generic.Emph:          '#008800',
        token.Generic.Error:         '#d30202',
        token.Error:                 'bg:#e3d2d2 #a61717'
    }

def format_run(ttype, run):
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    color, bold, italic, underline, _, border, roman, sans, mono = SasStyle._styles[ttype]
    
    if color:
        if len(color) == 3:
            color = color[0]*2, color[1]*2, color[2]*2
        else:
            color = color[:2], color[2:4], color[4:]
        color = [int(c, 16) for c in color]
        run.font.color.rgb = RGBColor(*color)

class PowerPointCodeFormatter(Formatter):
    MAX_REASONABLE_LINE = 14
    def __init__(self, text_frame, code, **options):
        self.options = options
        self.text_frame = text_frame
        self.code = code
        self.line_count = code.count('\n')
        # 9 fits comfortably

    def fix_height(self):
        paragraph = self.text_frame.paragraphs[0]
        #if self.line_count > 9:
        #    spacing = (self.MAX_REASONABLE_LINE-(self.line_count-9))/self.MAX_REASONABLE_LINE/2
        #    paragraph.line_spacing = spacing
        paragraph.line_spacing = .5

    def format(self, tokensource, outfile):
        if self.text_frame:
            self.text_frame.clear()
            paragraph = self.text_frame.paragraphs[0]
            no_bullet(paragraph)
            paragraph.font.name = "Courier New"
            for ttype, value in tokensource:
                run = paragraph.add_run()
                run.text = value
                format_run(ttype, run)
            self.fix_height()
        else:
            print("Throwing away codeblock!")


class PowerPointRenderer(GFMRendererMixin):
    options = {}
    # TODO: Fix these to be instance locals instead of class locals!
    narrate = False
    voice = polly.DEFAULT_VOICE
    GRAPHICS_FOLDER = "./"
    BASE_PRESENTATION = POWERPOINT_TEMPLATE
    SLIDE_LAYOUT_TYPES = {
        "title": 0,
        "title_content": 1,
        "section": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "captioned_content": 7,
        "captioned_picture": 8,
    }

    def __init__(self, **options):
        self.current = None
        self._current_text = None
        self._current_slide = None
        self.is_blank_slide = True
        self.presentation = Presentation(self.BASE_PRESENTATION)
        self._list = []
        self._notes = []
        self._durations = []
        self._transcript = []
        self._seen_summary = False

    @property
    def current_slide(self):
        if self._current_slide is None:
            self._current_slide = self.add_slide("blank")
            self._current_text = None
        return self._current_slide

    @property
    def current_text(self):
        if self._current_text is None:
            new_textbox = self.current_slide.shapes.add_textbox()
            self._current_text = new_textbox.text_frame
        return self._current_text

    def finish_previous_slides(self):
        if self._notes:
            notes = "\n".join(n for n in self._notes if n)
            self.add_narration(self._current_slide, notes)
            self._transcript.append(notes)
            self._notes = []

    def add_slide(self, type="title"):
        slide_layout_index = self.SLIDE_LAYOUT_TYPES.get(type, 6)
        slide_layout = self.presentation.slide_layouts[slide_layout_index]
        self._current_slide = self.presentation.slides.add_slide(slide_layout)
        if type == "title_content":
            self._current_text = self.current_slide.shapes[1].text_frame
        self.is_blank_slide = True
        return self._current_slide

    @staticmethod
    def autoplay_media(media):
        el_id = xpath(media.element, ".//p:cNvPr")[0].attrib["id"]
        el_cnt = xpath(
            media.element.getparent().getparent().getparent(),
            './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
        )[0]
        cond = xpath(el_cnt.getparent().getparent(), ".//p:cond")[0]
        cond.set("delay", "1000")

    def add_slide_transition(self, slide, duration):
        xpath(slide.element, ".//p:cSld")[0].addnext(
            etree.fromstring(
                f"""
            <mc:AlternateContent
                xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"
                xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
                xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main">
                <mc:Choice Requires="p159">
                    <p:transition spd="fast" p14:dur="1000" advTm="{duration}">
                        <p159:morph option="byObject" />
                    </p:transition>
                </mc:Choice>
                <mc:Fallback>
                    <p:transition spd="fast" advTm="{duration}">
                        <p:fade />
                    </p:transition>
                </mc:Fallback>
            </mc:AlternateContent>"""
            )
        )

    def add_audio_overlay(self, slide, audio_file) -> int:
        #print(audio_file)
        seconds = math.ceil(MP3(audio_file).info.length + 2)
        duration = seconds * 1000
        self.add_slide_transition(slide, duration)
        #audio_file = os.path.abspath(audio_file)
        audio_object = slide.shapes.add_movie(audio_file, left=Inches(0),
        top=Inches(0), width=Inches(0), height=Inches(0))#, poster_frame_image = None)
        self.autoplay_media(audio_object)
        self.is_blank_slide = False
        return seconds

    def add_narration(self, slide, text):
        audio_file = polly.speech(text, self.voice, self.narrate, label=self._input_path)
        duration = self.add_audio_overlay(slide, audio_file)
        self._durations.append(duration)

    def render_strong_emphasis(self, element: "inline.StrongEmphasis") -> str:
        return f"{self.render_children(element)}"

    def render_emphasis(self, element: "inline.Emphasis") -> str:
        return self.render_children(element)

    def render_code_span(self, element: "inline.CodeSpan") -> str:
        text = element.children
        if text and text[0] == "`" or text[-1] == "`":
            return f"{text}"
        return str(text)
        #return f"{element.children}"

    def render_fenced_code(self, element):
        code = element.children[0].children
        options = PowerPointRenderer.options.copy()
        # options.update(_parse_extras(getattr(element, "extra", None)))
        if element.lang:
            try:
                lexer = get_lexer_by_name(element.lang, stripall=True)
            except ClassNotFound:
                lexer = guess_lexer(code)
        else:
            lexer = guess_lexer(code)

        if code.count('\n') < PowerPointCodeFormatter.MAX_REASONABLE_LINE:
            formatter = PowerPointCodeFormatter(self.current_text, code, **options)
            result = highlight(code, lexer, formatter)
        else:
            formatter = ImageFormatter(line_numbers=False, style = CodeStyle,#get_style_by_name('sas'), 
                font_size=30, image_pad = 0, line_pad = 8, **options)
            image_information = highlight(code, lexer, formatter)
            with io.BytesIO() as temporary_image:
                temporary_image.write(image_information)
                temporary_image.seek(0)
                placeholder = self.current_slide.placeholders[1]
                replace_with_image(temporary_image, placeholder, self.current_slide, True, self.presentation)
            # For no real reason, also generate HTML
            formatter = HtmlFormatter(**options)
            result = highlight(code, lexer, formatter)
            
        #self.current_text.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        #self.current_text.fit_text("Courier New")
        return result
    
    def add_transcript(self, text):
        self._notes.append(text)

    def is_summary(self, element):
        return element.children and element.children[0].children and  element.children[0].children == "Summary"

    def render_heading(self, element: "block.Heading") -> str:
        self.finish_previous_slides()
        child_content = self.render_children(element)
        if self.is_summary(element):
            self._seen_summary = True
            if self._seen_summary:
                return child_content
        if element.level == 1:
            new_slide = self.add_slide("title")
            self.add_transcript(child_content)
        else:
            new_slide = self.add_slide("title_content")
        new_slide.shapes.title.text = child_content
        if element.level == 1:
            new_slide.placeholders[1].text = "The Python Bakery"
        return "<h{level}>{children}</h{level}>\n".format(
            level=element.level, children=child_content
        )

    def render_paragraph(self, element: "block.Paragraph") -> str:
        children = self.render_children(element)
        if self._list or self._seen_summary:
            return children
        if self.current_slide.has_notes_slide:
            self.current_slide.notes_slide.notes_text_frame.text += "\n"
        self.current_slide.notes_slide.notes_text_frame.text += children
        self.add_transcript(children)
        self.is_blank_slide = False
        if element._tight:  # type: ignore
            return children
        else:
            return f"<p>{children}</p>\n"

    def render_list(self, element: "block.List") -> str:
        if self._seen_summary:
            return ""
        self._list.append(element)
        children = self.render_children(element)
        # PowerPoint output
        self.current_text.text += children
        # Regular markdown output
        if element.ordered:
            tag = "ol"
            extra = f' start="{element.start}"' if element.start != 1 else ""
        else:
            tag = "ul"
            extra = ""
        self._list.pop()
        return ""
        # return "<{tag}{extra}>\n{children}</{tag}>\n".format(
        #    tag=tag, extra=extra, children=children
        # )

    def render_list_item(self, element: "block.ListItem") -> str:
        children = self.render_children(element)
        if len(element.children) == 1 and getattr(element.children[0], "_tight", False):  # type: ignore
            sep = ""
        else:
            sep = "\n"
        return f"{sep}{children}\n"
        # return f"<li>{sep}{children}</li>\n"

    def render_image(self, element: "inline.Image") -> str:
        url = self.escape_url(os.path.join(self.GRAPHICS_FOLDER, element.dest))
        # for shape in self.current_slide.placeholders:
        #    print('%d %s %s' % (shape.placeholder_format.idx, shape.name, shape.placeholder_format.type), dir(shape))
        placeholder = self.current_slide.placeholders[1]
        replace_with_image(url, placeholder, self.current_slide)
        render_func = self.render
        self.render = self.render_plain_text  # type: ignore
        body = self.render_children(element)
        self.render = render_func  # type: ignore
        #return template.format(url, body, title)
        return ""

    def finish(self):
        self.finish_previous_slides()
        return ""

    @staticmethod
    def escape_html(raw: str) -> str:
        return raw

    @staticmethod
    def escape_url(raw: str) -> str:
        return raw


class PPTXRenderExtension:
    elements = [
        elements.Paragraph,
        elements.ListItem,
        elements.Strikethrough,
        elements.Url,
        elements.Table,
        elements.TableRow,
        elements.TableCell,
    ]
    renderer_mixins = [PowerPointRenderer]
    parser_mixins = []


# XML Stuff

ETREE_NAMESPACE_MAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "p159": "http://schemas.microsoft.com/office/powerpoint/2015/09/main"
}

for k, v in ETREE_NAMESPACE_MAP.items():
    etree.register_namespace(k, v)


def xpath(el, query):
    return etree.ElementBase.xpath(el, query, namespaces=ETREE_NAMESPACE_MAP)

CREATE_VIDEO_STATUSES = {
    0: "None",
    1: "In Progress",
    2: "In Queue",
    3: "Done",
    4: "Failed"
}

# https://docs.microsoft.com/en-us/office/vba/api/powerpoint.presentation.createvideo
ppSaveAsWMV, ppSaveAsMP4 = 37, 39
def convert_ppt_to_wmv(ppt_src, wmv_target, fps=24, quality=100, resolution=1080):
    ppt_src, wmv_target = os.path.abspath(ppt_src), os.path.abspath(wmv_target)
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    presentation = ppt.Presentations.Open(ppt_src, WithWindow=False)
    presentation.CreateVideo(wmv_target,-1,4,resolution,fps,quality)
    start_time_stamp = time.time()
    print(f"Status: {CREATE_VIDEO_STATUSES[presentation.CreateVideoStatus]}")
    with tqdm() as pbar:
        while presentation.CreateVideoStatus < 3:
            pbar.update(n=1)
            pbar.set_description(f"Status: {CREATE_VIDEO_STATUSES[presentation.CreateVideoStatus]}")
            #pbar.refresh()
            time.sleep(1)
    end_time_stamp=time.time()
    print(end_time_stamp-start_time_stamp)
    time.sleep(1)
    ppt.Quit()


def convert_wmv_to_mp4(wmv_src, mp4_target):
    wmv_src, mp4_target = os.path.abspath(wmv_src), os.path.abspath(mp4_target)
    ffmpeg.input(wmv_src).output(mp4_target).run(overwrite_output=True, quiet=True)


HASH_FILE_PATH = "finished_video_hashes.json"

def load_previous_hashes():
    # Create the hash file if it doesn't exist
    if not os.path.exists(HASH_FILE_PATH):
        with open(HASH_FILE_PATH, 'w') as hash_file:
            json.dump({}, hash_file)
    # Actually load it
    with open(HASH_FILE_PATH) as hash_file:
        return json.load(hash_file)

def save_hash(contents):
    hashes = load_previous_hashes()
    hashes[hash(contents)] = contents
    with open(HASH_FILE_PATH, 'w') as hash_file:
        json.dump(hashes, hash_file, indent=2)

# Main Function

WMV_OPTIONS = {
    'low': {'quality': 40, 'resolution': 720},
    'high': {'quality': 100, 'resolution': 1080}
}

def bake_markdown(input_path, output_path, graphics_path, narrate, voice, wmv, force_rebuild, nosave, transcript, mp4):
    previous_hashes = load_previous_hashes()
    PowerPointRenderer.GRAPHICS_FOLDER = graphics_path
    PowerPointRenderer.narrate = narrate
    PowerPointRenderer.voice = voice
    PowerPointRenderer._input_path = input_path
    converter = marko.Markdown()
    converter.use(PPTXRenderExtension)
    with open(input_path, encoding='utf-8') as input_file:
        input_text = input_file.read()
    if output_path is None:
        output_path = Path(input_path).stem
        if output_path.endswith('_read'):
            output_path = output_path[:-len('_read')]
        output_path = os.path.join('../build/', output_path)
    regular_metadata, front_matter_metadata, input_content = extract_front_matter(input_text)
    if not force_rebuild and hash_exists(input_text, previous_hashes):
        yield "Skipping - hashed output already exists: " + previous_hashes[hash(input_text)]
    else:
        rendered = converter.convert(input_content)
        rendered += converter.renderer.finish()
        if nosave:
            yield "Skipping - nosave parameter was given."
        else:
            with open(output_path + ".html", "w", encoding='utf-8') as output_file:
                output_file.write(rendered)
            presentation = converter.renderer.presentation
            presentation.save(output_path + f"-{voice}.pptx")
            yield "Finished powerpoint"
            if wmv != 'none':
                convert_ppt_to_wmv(output_path+f"-{voice}.pptx", output_path+f"-{voice}.wmv", **WMV_OPTIONS[wmv])
                save_hash(input_text)
                yield "Finished wmv"
            if mp4:
                convert_wmv_to_mp4(output_path+f"-{voice}.wmv", output_path+f"-{voice}.mp4")
                yield "Finished mp4"
            if transcript:
                with open(f"{output_path}-{voice}.vtt", "w", encoding='utf-8') as output_file:
                    captions = make_captions(converter.renderer._transcript, converter.renderer._durations)
                    output_file.write("\n".join(captions))
                    yield "Finished captions"


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Compile Markdown into PowerPoint Videos"
    )
    parser.add_argument("input", metavar="i", help="The input Markdown file (.md)")
    parser.add_argument(
        "--output", metavar="o",
        help="The base filename for the outputs (e.g., PowerPoint file, WMV file). If not provided, then the path will be generated based on the input filename.",
        default=None
    )
    parser.add_argument("--graphics", metavar="g", help="The location of the folder with images in it.", default="../graphics/")

    parser.add_argument('-a', "--narrate", action='store_true', help="Add in automatic narration using Amazon Polly.")
    parser.add_argument('-v', "--voice", choices=['Amy', 'Bart'], default=polly.DEFAULT_VOICE, help="Choose the voice-over files that will be used.")

    parser.add_argument("-w", "--wmv", choices=['none', 'low', 'high'], default='none', help="Export a WMV file too")
    parser.add_argument("-m", "--mp4", action="store_true", help="Export an MP4 file too")

    parser.add_argument("-f", "--force", action="store_true", help="Force recreating the built output, even if its current hash is in the built list.")
    parser.add_argument("-n", "--nosave", action="store_true",
                        help="Do NOT save the rendered presentation at all. Useful for setting up narration, since audio files will still be created.")
    
    parser.add_argument('-t', "--transcript", action="store_true", help="Generate a transcript of the narration.")

    args = parser.parse_args()
    for progress in bake_markdown(args.input, args.output, args.graphics, args.narrate, args.voice,
                                    args.wmv, args.force, args.nosave, args.transcript, args.mp4):
        print(progress)
